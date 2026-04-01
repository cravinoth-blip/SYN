"""
PPTX Comparator Router
Endpoints for parsing, comparing, and editing PowerPoint files.
Mount in app.py with: app.include_router(router)
"""
import io
import json
import os
import logging
from typing import Any

from fastapi import APIRouter, Form, HTTPException, UploadFile
from fastapi.responses import StreamingResponse
from openai import OpenAI
from pptx import Presentation
from pptx.util import Pt

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/pptx", tags=["pptx"])


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _extract_slide_shapes(slide) -> list[dict[str, Any]]:
    """Return a list of shape dicts with id, name, and text for a single slide."""
    shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = "\n".join(
                para.text for para in shape.text_frame.paragraphs
            )
            shapes.append({
                "shape_id": shape.shape_id,
                "shape_name": shape.name,
                "text": text,
            })
    return shapes


def _extract_slide_text(slide) -> str:
    """Return all text from a slide as a single string."""
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
    return "\n".join(parts)


def _load_pptx(upload: UploadFile) -> Presentation:
    """Read an UploadFile into a python-pptx Presentation."""
    data = upload.file.read()
    return Presentation(io.BytesIO(data))


def _apply_edits(prs: Presentation, edits: list[dict[str, Any]]) -> None:
    """
    Apply a list of edits to a Presentation in-place.
    Each edit: {"slide_idx": int, "shape_name": str, "new_text": str}
    Matches shapes by name, replaces text by clearing runs and setting first run.
    """
    for edit in edits:
        slide_idx = int(edit.get("slide_idx", 0))
        shape_name = edit.get("shape_name", "")
        new_text = edit.get("new_text", "")

        if slide_idx < 0 or slide_idx >= len(prs.slides):
            logger.warning("Edit references out-of-range slide_idx=%d; skipping.", slide_idx)
            continue

        slide = prs.slides[slide_idx]
        matched = False
        for shape in slide.shapes:
            if shape.name == shape_name and shape.has_text_frame:
                tf = shape.text_frame
                # Work only with the first paragraph's runs for simplicity;
                # preserve paragraph-level formatting by clearing extra paragraphs.
                for para in tf.paragraphs:
                    if not para.runs:
                        continue
                    # Keep first run, set its text; remove remaining runs.
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run.text = ""
                    # After writing the first paragraph, blank any subsequent ones.
                    new_text = ""
                matched = True
                break

        if not matched:
            logger.warning(
                "Edit: shape '%s' not found on slide %d; skipping.",
                shape_name,
                slide_idx,
            )


# ---------------------------------------------------------------------------
# Endpoints
# ---------------------------------------------------------------------------

@router.post("/slides")
async def parse_slides(file_a: UploadFile, file_b: UploadFile):
    """
    Parse two PPTX files and return their slide/shape structure.
    """
    try:
        prs_a = _load_pptx(file_a)
        prs_b = _load_pptx(file_b)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Failed to parse PPTX: {exc}") from exc

    slides_a = [
        {"slide_idx": i, "shapes": _extract_slide_shapes(slide)}
        for i, slide in enumerate(prs_a.slides)
    ]
    slides_b = [
        {"slide_idx": i, "shapes": _extract_slide_shapes(slide)}
        for i, slide in enumerate(prs_b.slides)
    ]

    return {
        "slide_count_a": len(prs_a.slides),
        "slide_count_b": len(prs_b.slides),
        "slides_a": slides_a,
        "slides_b": slides_b,
    }


@router.post("/compare-slide")
async def compare_slide(
    file_a: UploadFile,
    file_b: UploadFile,
    slide_idx: int = Form(0),
    model: str = Form("gpt-4o"),
):
    """
    Extract text from a single slide in both files and use OpenAI to compare them.
    Returns structured JSON with differences, suggestions, and an overall assessment.
    """
    try:
        prs_a = _load_pptx(file_a)
        prs_b = _load_pptx(file_b)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Failed to parse PPTX: {exc}") from exc

    if slide_idx < 0 or slide_idx >= len(prs_a.slides):
        raise HTTPException(
            status_code=400,
            detail=f"slide_idx {slide_idx} is out of range for Presentation A ({len(prs_a.slides)} slides).",
        )
    if slide_idx >= len(prs_b.slides):
        raise HTTPException(
            status_code=400,
            detail=f"slide_idx {slide_idx} is out of range for Presentation B ({len(prs_b.slides)} slides).",
        )

    text_a = _extract_slide_text(prs_a.slides[slide_idx])
    text_b = _extract_slide_text(prs_b.slides[slide_idx])

    api_key = os.getenv("OPENAI_API_KEY", "")
    if not api_key:
        raise HTTPException(status_code=500, detail="OPENAI_API_KEY is not configured.")

    client = OpenAI(api_key=api_key)

    system_prompt = (
        "You are an expert presentation analyst. "
        "Compare two PowerPoint slides and return a JSON object with exactly these keys:\n"
        "- summary: string — a concise overall summary of the comparison.\n"
        "- differences: array of objects with keys aspect (string), slide_a (string), slide_b (string).\n"
        "- suggestions: array of objects with keys shape_name (string), original_text (string), "
        "  suggested_text (string), reason (string). These are suggested edits for Slide A.\n"
        "- overall_assessment: string — your final recommendation.\n"
        "Return only valid JSON, no markdown."
    )
    user_prompt = (
        f"=== SLIDE A (slide index {slide_idx}) ===\n{text_a}\n\n"
        f"=== SLIDE B (slide index {slide_idx}) ===\n{text_b}"
    )

    try:
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            response_format={"type": "json_object"},
        )
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"OpenAI API error: {exc}") from exc

    raw = response.choices[0].message.content or "{}"
    try:
        result = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise HTTPException(status_code=502, detail=f"Failed to parse OpenAI JSON response: {exc}") from exc

    # Ensure required keys exist with sane defaults
    result.setdefault("summary", "")
    result.setdefault("differences", [])
    result.setdefault("suggestions", [])
    result.setdefault("overall_assessment", "")

    return result


@router.post("/edit")
async def edit_pptx(
    file_a: UploadFile,
    edits: str = Form(...),
):
    """
    Apply a list of text edits to Presentation A and return the modified file.
    edits: JSON string — list of {"slide_idx": int, "shape_name": str, "new_text": str}
    """
    try:
        prs = _load_pptx(file_a)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Failed to parse PPTX: {exc}") from exc

    try:
        edit_list: list[dict[str, Any]] = json.loads(edits)
    except json.JSONDecodeError as exc:
        raise HTTPException(status_code=422, detail=f"Invalid edits JSON: {exc}") from exc

    if not isinstance(edit_list, list):
        raise HTTPException(status_code=422, detail="edits must be a JSON array.")

    _apply_edits(prs, edit_list)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = (file_a.filename or "presentation").replace(".pptx", "") + "_edited.pptx"

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
