import io
from dataclasses import dataclass

from pptx import Presentation
from pypdf import PdfReader


@dataclass(frozen=True)
class ParsedContent:
    text: str
    metadata: dict


def parse_pdf(data: bytes) -> ParsedContent:
    reader = PdfReader(io.BytesIO(data))
    pages = []
    for index, page in enumerate(reader.pages):
        pages.append({"page": index + 1, "text": page.extract_text() or ""})
    text = "\n\n".join(page["text"] for page in pages if page["text"])
    return ParsedContent(text=text, metadata={"page_count": len(reader.pages), "parser": "pypdf"})


def parse_pptx(data: bytes) -> ParsedContent:
    presentation = Presentation(io.BytesIO(data))
    slide_text: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        if parts:
            slide_text.append(f"Slide {slide_index}\n" + "\n".join(parts))
    return ParsedContent(
        text="\n\n".join(slide_text),
        metadata={"slide_count": len(presentation.slides), "parser": "python-pptx"},
    )


def parse_document(filename: str, data: bytes) -> ParsedContent:
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return parse_pdf(data)
    if lower.endswith(".pptx"):
        return parse_pptx(data)
    raise ValueError("Only PDF and PPTX files are supported")


def chunk_text(text: str, chunk_size: int = 1800, overlap: int = 160) -> list[str]:
    clean = " ".join(text.split())
    if not clean:
        return []
    chunks: list[str] = []
    cursor = 0
    while cursor < len(clean):
        chunks.append(clean[cursor : cursor + chunk_size])
        cursor += chunk_size - overlap
    return chunks
