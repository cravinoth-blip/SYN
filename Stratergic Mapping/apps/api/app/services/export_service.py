import io
import uuid

from fastapi import HTTPException
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from sqlalchemy.orm import Session

from app.models import ExportJob, ExportType, JobStatus, Version
from app.services.audit import record_audit
from app.services.storage import storage
from app.services.workspace_service import build_workspace


def create_export_job(db: Session, version_id: uuid.UUID, export_type: ExportType) -> ExportJob:
    version = db.get(Version, version_id)
    if not version:
        raise HTTPException(status_code=404, detail="Version not found")
    job = ExportJob(version_id=version_id, export_type=export_type, status=JobStatus.RUNNING)
    db.add(job)
    db.flush()
    try:
        workspace = build_workspace(db, version.project_id, version_id=version_id)
        if export_type == ExportType.PDF:
            data = _render_pdf(workspace)
            extension = "pdf"
            content_type = "application/pdf"
        else:
            data = _render_pptx(workspace)
            extension = "pptx"
            content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        key = f"projects/{version.project_id}/exports/{job.export_job_id}.{extension}"
        job.artifact_uri = storage.put_bytes(key, data, content_type=content_type)
        job.status = JobStatus.SUCCEEDED
        record_audit(
            db,
            "export.completed",
            project_id=version.project_id,
            version_id=version_id,
            payload={"export_job_id": str(job.export_job_id), "export_type": export_type.value},
        )
    except Exception as exc:
        job.status = JobStatus.FAILED
        record_audit(
            db,
            "export.failed",
            project_id=version.project_id,
            version_id=version_id,
            payload={"export_job_id": str(job.export_job_id), "error": str(exc)},
        )
    db.commit()
    db.refresh(job)
    return job


def get_export_job(db: Session, export_job_id: uuid.UUID) -> ExportJob:
    job = db.get(ExportJob, export_job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Export job not found")
    return job


def _render_pdf(workspace) -> bytes:
    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    y = height - 48
    pdf.setTitle(f"{workspace.project.project_name} 7Cs Report")
    pdf.setFont("Helvetica-Bold", 16)
    pdf.drawString(48, y, workspace.project.project_name[:80])
    y -= 28
    pdf.setFont("Helvetica", 10)
    pdf.drawString(48, y, f"{workspace.project.disease} | {workspace.project.geography}")
    y -= 28
    for section in workspace.sections:
        if y < 120:
            pdf.showPage()
            y = height - 48
        pdf.setFont("Helvetica-Bold", 13)
        pdf.drawString(48, y, section.section_name)
        y -= 18
        pdf.setFont("Helvetica", 9)
        for line in section.narrative_markdown.replace("#", "").splitlines():
            for part in _wrap(line, 95):
                if y < 64:
                    pdf.showPage()
                    y = height - 48
                    pdf.setFont("Helvetica", 9)
                pdf.drawString(48, y, part)
                y -= 12
        y -= 14
    pdf.showPage()
    pdf.save()
    return buffer.getvalue()


def _render_pptx(workspace) -> bytes:
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = workspace.project.project_name
    title_slide.placeholders[1].text = f"{workspace.project.disease} | {workspace.project.geography}"
    for section in workspace.sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section.section_name
        body = slide.placeholders[1].text_frame
        body.text = section.narrative_markdown.replace("#", "")[:1200]
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _wrap(text: str, width: int) -> list[str]:
    if not text:
        return [""]
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if len(candidate) > width:
            lines.append(current)
            current = word
        else:
            current = candidate
    if current:
        lines.append(current)
    return lines

